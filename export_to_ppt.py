"""
Export the ChiCom dashboard to a PowerPoint deck.

Layout rules:
  - Each chart on the dashboard is one individual image in screenshots/.
  - "Square-ish" charts (aspect <= 2.0) are placed 2-per-slide side-by-side.
  - "Wide" charts (aspect > 2.0, like col-12 trend lines) get their own slide,
    centered full-width.
  - Q sections with >2 charts split across multiple slides ("(tiếp)" suffix).
  - Insight text box (editable) appears on the FIRST slide of each Q, under
    the charts. Subsequent slides of the same Q omit it.

Prerequisites:
    python build_data.py          # produces data_computed.js + insights
    python screenshot_dashboard.py # produces screenshots/*.png + manifest.json

Run:
    python export_to_ppt.py
    python export_to_ppt.py --out chicom_report.pptx
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = Path(__file__).resolve().parent
SCREENSHOT_DIR = BASE / "screenshots"
MANIFEST_PATH = SCREENSHOT_DIR / "manifest.json"


# ── Style ────────────────────────────────────────────────────────────────

ACCENT     = RGBColor(0x3A, 0x56, 0xC9)
PURPLE     = RGBColor(0x6E, 0x41, 0xA6)
DARK_TEXT  = RGBColor(0x1C, 0x20, 0x2C)
MUTED_TEXT = RGBColor(0x66, 0x6C, 0x78)
GRAY       = RGBColor(0x80, 0x80, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout constants (inches)
HEADER_TOP      = 0.24
HEADER_H        = 0.60
SUBTITLE_H      = 0.30
CHART_TOP       = 1.25   # just below the subtitle
FOOTER_TOP      = 7.18

INSIGHT_TOP     = 5.70   # when insight is present
INSIGHT_H       = 1.40
CHART_H_W_INS   = 4.30   # chart area height when insight present
CHART_H_NO_INS  = 5.80   # chart area height when no insight


# ── Text helpers ─────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_para(p, text: str, *, size: int, bold: bool = False,
              color: RGBColor = DARK_TEXT, font: str = "Inter"):
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _header_bar(slide, title: str, subtitle: str = ""):
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.16)
    )
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(HEADER_TOP), Inches(12.5), Inches(HEADER_H)
    )
    title_box.text_frame.word_wrap = True
    title_box.text_frame.margin_left = Emu(0)
    _set_para(title_box.text_frame.paragraphs[0], title, size=20, bold=True)

    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.5), Inches(HEADER_TOP + HEADER_H + 0.02),
            Inches(12.5), Inches(SUBTITLE_H)
        )
        sub.text_frame.margin_left = Emu(0)
        _set_para(sub.text_frame.paragraphs[0], subtitle, size=11, color=MUTED_TEXT)


def _footer(slide, page: int, total: int, date_range: dict):
    start = (date_range or {}).get("start") or "—"
    end = (date_range or {}).get("end") or "—"
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(FOOTER_TOP), Inches(12.5), Inches(0.3))
    foot.text_frame.margin_left = Emu(0)
    _set_para(
        foot.text_frame.paragraphs[0],
        f"ChiCom · Community Insights  ·  {start} → {end}  ·  trang {page}/{total}",
        size=9, color=GRAY,
    )


def _insight_box(slide, paragraph: str, *, top: float = INSIGHT_TOP,
                 height: float = INSIGHT_H,
                 label: str = "AI INSIGHT — chỉnh sửa tự do"):
    if not paragraph:
        return
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(top), Inches(0.08), Inches(height)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE

    body = slide.shapes.add_textbox(
        Inches(0.7), Inches(top), Inches(12.2), Inches(height)
    )
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(30000)

    _set_para(tf.paragraphs[0], label, size=9, bold=True, color=PURPLE)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    _set_para(p, paragraph, size=12, color=DARK_TEXT)


# ── Image placement ──────────────────────────────────────────────────────

def _fit_into(img_path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    with Image.open(img_path) as im:
        w, h = im.size
    aspect = w / h if h else 1
    if (max_w / max_h) > aspect:
        fit_h = max_h
        fit_w = fit_h * aspect
    else:
        fit_w = max_w
        fit_h = fit_w / aspect
    return fit_w, fit_h


def _add_image(slide, img_path: Path, *, left: float, top: float,
               max_w: float, max_h: float):
    fit_w, fit_h = _fit_into(img_path, max_w, max_h)
    # Center within the allotted box
    x = left + (max_w - fit_w) / 2
    y = top + (max_h - fit_h) / 2
    pic = slide.shapes.add_picture(
        str(img_path),
        Inches(x), Inches(y),
        width=Inches(fit_w), height=Inches(fit_h),
    )
    pic.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    pic.line.width = Pt(0.5)
    return pic


# ── Layout planning ──────────────────────────────────────────────────────

WIDE_ASPECT_THRESHOLD = 2.0  # aspect > this → 1 chart per slide


def _plan_slides(images: list[dict]) -> list[list[dict]]:
    """
    Group images into slide-sized batches.
      wide (aspect > threshold) → solo on slide
      square (aspect <= threshold) → paired 2-per-slide
    Preserves original order.
    """
    slides: list[list[dict]] = []
    pending_square: dict | None = None

    for img in images:
        if img["aspect"] > WIDE_ASPECT_THRESHOLD:
            # flush pending square as solo
            if pending_square is not None:
                slides.append([pending_square])
                pending_square = None
            slides.append([img])  # wide solo
        else:
            if pending_square is None:
                pending_square = img
            else:
                slides.append([pending_square, img])
                pending_square = None

    if pending_square is not None:
        slides.append([pending_square])  # trailing solo square
    return slides


# ── Data capture ─────────────────────────────────────────────────────────

def _load_manifest() -> dict:
    if not MANIFEST_PATH.is_file():
        sys.exit(
            f"ERROR: {MANIFEST_PATH} not found. Run `python screenshot_dashboard.py` first."
        )
    return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))


def _compute_info_and_date_range() -> tuple[dict, dict]:
    sys.path.insert(0, str(BASE / "backend"))
    import re
    import pandas as pd
    from compute import compute_all

    csv_path = BASE / "output" / "filtered" / "all_groups_final_v2_annotated.csv"
    if not csv_path.is_file():
        sys.exit(f"ERROR: CSV not found at {csv_path}.")
    df = pd.read_csv(csv_path, encoding="utf-8-sig", low_memory=False)
    js, info = compute_all(df)

    def _extract(name: str) -> dict | None:
        m = re.search(rf"const {name}\s*=\s*(\{{[\s\S]*?\n\s*\}});", js)
        if not m:
            return None
        try:
            return json.loads(m.group(1))
        except Exception:
            return None

    return info, _extract("DATE_RANGE") or {}


def _load_insights() -> dict[str, str]:
    out: dict[str, str] = {}
    manual = BASE / "backend" / "insights_manual.json"
    if manual.is_file():
        data = json.loads(manual.read_text(encoding="utf-8"))
        for k, v in data.items():
            if k.startswith("Q") and isinstance(v, str) and v.strip():
                out[k] = v.strip()
    cache = BASE / ".insight_cache.json"
    if cache.is_file():
        try:
            c = json.loads(cache.read_text(encoding="utf-8"))
            for k, entry in c.items():
                if entry and entry.get("insight"):
                    out[k] = entry["insight"]
        except Exception:
            pass
    return out


# ── High-level slide builders ─────────────────────────────────────────────

def slide_title(prs: Presentation, info: dict, date_range: dict):
    slide = _blank(prs)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.3)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT

    title = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(1.0))
    _set_para(title.text_frame.paragraphs[0],
              "ChiCom — Community Insights Report",
              size=40, bold=True)

    subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(0.6))
    relevant = info.get("relevantPosts") or 0
    _set_para(
        subtitle.text_frame.paragraphs[0],
        f"Phân tích {relevant:,} Lượt Thảo Luận từ 9 cộng đồng seller Việt Nam · 14 câu hỏi nghiên cứu",
        size=16, color=MUTED_TEXT,
    )

    meta = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(12), Inches(0.5))
    start = (date_range or {}).get("start") or "—"
    end = (date_range or {}).get("end") or "—"
    months = (date_range or {}).get("monthsCount") or 0
    _set_para(
        meta.text_frame.paragraphs[0],
        f"Khoảng thời gian: {start} → {end}  ·  {months} tháng  ·  2 nhóm SOA + 7 nhóm EC",
        size=12, color=GRAY,
    )


def slide_section_banner(prs: Presentation, title: str, subtitle: str = ""):
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFC)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT

    t = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(12), Inches(1))
    _set_para(t.text_frame.paragraphs[0], title, size=34, bold=True)

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(12), Inches(0.5))
        _set_para(s.text_frame.paragraphs[0], subtitle, size=14, color=MUTED_TEXT)


def place_charts_on_slide(slide, batch: list[dict], *, has_insight: bool):
    """Place 1 or 2 images on the slide's chart area."""
    chart_h = CHART_H_W_INS if has_insight else CHART_H_NO_INS
    images = [SCREENSHOT_DIR / img["file"] for img in batch]

    if len(batch) == 1:
        # Solo chart — centered, up to 80% slide width
        _add_image(slide, images[0],
                   left=1.0, top=CHART_TOP,
                   max_w=11.33, max_h=chart_h)
        return

    # Two charts side-by-side — 6 in wide each, small gutter
    gutter = 0.3
    col_w = (12.33 - gutter) / 2
    _add_image(slide, images[0],
               left=0.5, top=CHART_TOP,
               max_w=col_w, max_h=chart_h)
    _add_image(slide, images[1],
               left=0.5 + col_w + gutter, top=CHART_TOP,
               max_w=col_w, max_h=chart_h)


# ── Q-section deck builder ───────────────────────────────────────────────

def build_q_slides(prs: Presentation, q_id: str, header_title: str,
                   subtitle: str, manifest: dict, insight: str | None):
    """Create all slides for one Q section (or the overview)."""
    entries = manifest.get(q_id) or []
    if not entries:
        # Still emit a slide with just the insight so the Q isn't missing
        slide = _blank(prs)
        _header_bar(slide, header_title, subtitle)
        if insight:
            _insight_box(slide, insight, top=CHART_TOP + 0.3,
                         height=SLIDE_H.inches - CHART_TOP - 0.8)
        return

    slide_batches = _plan_slides(entries)
    for idx, batch in enumerate(slide_batches):
        slide = _blank(prs)
        suffix = "" if idx == 0 else "  (tiếp)"
        _header_bar(slide, f"{header_title}{suffix}", subtitle if idx == 0 else "")
        has_insight = (idx == 0) and bool(insight)
        place_charts_on_slide(slide, batch, has_insight=has_insight)
        if has_insight:
            _insight_box(slide, insight)


# ── Top-level deck ───────────────────────────────────────────────────────

SECTIONS: list[tuple[str, str, str]] = [
    # (section_id, header title, subtitle)
    ("overview", "Phần 1 — Thông tin sơ bộ",
     "SOV community · Phân bố persona · Tổng quan dataset"),
    ("Q1", "Q1 — Master Topics: tỉ trọng tổng thể",
     "% đề cập trung bình qua các nhóm + phân bố theo SOA / EC"),
    ("Q2", "Q2 — Master Topics × Persona",
     "Heatmap cross-tab: nhóm nào thảo luận topic nào"),
    ("Q3", "Q3 — Seller vs Prospect",
     "Khác biệt Master Topics + sub-topic diverging"),
    ("Q4", "Q4 — Xu hướng Master Topics theo tháng",
     "Monthly · stacked % · weekly với tuần đỉnh"),
    ("Q5", "Q5 / Q6 — Chủ đề tiêu cực theo thời gian",
     "Ngày × giờ · heatmap · khung giờ cao điểm auto-detect"),
    ("Q7", "Q7 — Lý do gia nhập Amazon + benefit",
     "Top triggers · top benefits · sentiment split"),
    ("Q8", "Q8 — Dấu hiệu rời bỏ Amazon (SOA)",
     "Top lý do · persona rời bỏ · monthly trend"),
    ("Q9", "Q9 — Nhóm tham gia Q7 vs Q8",
     "Persona donuts cho từng luồng"),
    ("Q10", "Q10 — Top ngành hàng được thảo luận",
     "Top 10 categories · weekly trend"),
    ("Q11", "Q11 — Mức sử dụng tool Amazon (SOA)",
     "Tool usage + hài lòng vs vấn đề + top factors"),
    ("Q12", "Q12 — Dịch vụ bên thứ 3 seller cần (SOA)",
     "Mentions vs Need · demand % · satisfaction"),
    ("Q13", "Q13 — Khóa học seller quan tâm (SOA)",
     "Mentions · seeking · interest · sentiment"),
    ("Q14", "Q14 — Chủ đề tăng trưởng + P&L (SOA)",
     "Distribution · top topics · sentiment breakdown"),
]


def build_deck(out_path: Path) -> None:
    manifest = _load_manifest()
    print("[ppt] capturing metadata…")
    info, date_range = _compute_info_and_date_range()
    insights = _load_insights()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("[ppt] assembling slides…")
    slide_title(prs, info, date_range)
    slide_section_banner(
        prs, "Phần 1 — Thông tin sơ bộ",
        "Community volume · persona split · baseline KPIs"
    )

    # Overview (Section 1 content)
    build_q_slides(prs, "overview",
                   "Phần 1 — Thông tin sơ bộ",
                   "SOV community + phân bố persona",
                   manifest, None)

    slide_section_banner(
        prs, "Phần 2 — Thông tin chi tiết",
        "14 câu hỏi nghiên cứu — charts từ dashboard + AI insight editable"
    )

    for sec_id, header, subtitle in SECTIONS:
        if sec_id == "overview":
            continue  # already handled above
        q_id = sec_id if sec_id != "Q5" else "Q5"  # Q5 covers Q5+Q6 charts
        build_q_slides(prs, sec_id, header, subtitle, manifest,
                       insights.get(q_id))

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        _footer(slide, i, total, date_range)

    prs.save(out_path)
    size_kb = out_path.stat().st_size / 1024
    print(f"[ppt] wrote {out_path}  ({size_kb:.1f} KB, {total} slides)")


if __name__ == "__main__":
    p = argparse.ArgumentParser()
    p.add_argument("--out", default="ChiCom_Report.pptx")
    args = p.parse_args()
    out = Path(args.out)
    if not out.is_absolute():
        out = BASE / out
    build_deck(out)
